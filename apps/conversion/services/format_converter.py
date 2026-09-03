import io
import csv
import os
import logging
from pdf2docx import Converter
import pandas as pd
import tabula
from PIL import Image
import fitz
from pptx import Presentation
def pdf_to_docx(pdf_bytes: bytes) -> bytes:
    input_stream = io.BytesIO(pdf_bytes)
    output_stream = io.BytesIO()
    cv = Converter(input_stream)
    cv.convert(output_stream, start=0, end=None)
    cv.close()
    output_stream.seek(0)
    return output_stream.read()

def pdf_to_xlsx(pdf_bytes: bytes) -> bytes:
    dfs = tabula.read_pdf(io.BytesIO(pdf_bytes), pages='all', multiple_tables=True)
    if not dfs:
        return b''
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        for i, df in enumerate(dfs):
            df.to_excel(writer, sheet_name=f'Table_{i+1}', index=False)
    output.seek(0)
    return output.read()

def pdf_to_csv(pdf_bytes: bytes) -> bytes:
    dfs = tabula.read_pdf(io.BytesIO(pdf_bytes), pages='all', multiple_tables=True)
    if not dfs:
        return b''
    output = io.StringIO()
     
    for i, df in enumerate(dfs):
        if i > 0:
            output.write('\n\n')
        df.to_csv(output, index=False)
    return output.getvalue().encode('utf-8')

def pdf_to_txt(pdf_bytes: bytes) -> bytes:
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
    return text.encode('utf-8')

def pdf_to_jpg(pdf_bytes: bytes, dpi: int = 150) -> bytes:
    # Convert first page to JPG
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    page = doc[0]
    pix = page.get_pixmap(dpi=dpi)
    img_bytes = pix.tobytes("jpeg")
    return img_bytes
def pdf_to_png(pdf_bytes: bytes, dpi: int = 150) -> bytes:
     
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    page = doc[0]
    pix = page.get_pixmap(dpi=dpi)
    return pix.tobytes("png")

def pdf_to_webp(pdf_bytes: bytes, dpi: int = 150, quality: int = 80) -> bytes:
     
    from PIL import Image
    import io
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    page = doc[0]
    pix = page.get_pixmap(dpi=dpi)
     
    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
    output = io.BytesIO()
    img.save(output, format="WEBP", quality=quality)
    return output.getvalue()



def pdf_to_pptx(pdf_bytes: bytes) -> bytes:

    try:
        
        prs = Presentation()
        
        
        prs.slide_width = 9144000  
        prs.slide_height = 5143500  
        
        # Open PDF
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        for page_num in range(len(doc)):
            # Render page to image
            page = doc[page_num]
            pix = page.get_pixmap(dpi=150)
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))
            
            # Create blank slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])  
            
            # Add image as background
            img_stream = io.BytesIO()
            img.save(img_stream, format='PNG')
            img_stream.seek(0)
            
            slide.shapes.add_picture(
                img_stream,
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height
            )
        
        doc.close()
        
        # Save presentation
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.getvalue()
        
    except Exception as e:
        logging.exception(f"Failed to convert PDF to PPTX: {e}")
        raise
    
    
    